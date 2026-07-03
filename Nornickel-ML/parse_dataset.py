from __future__ import annotations

import argparse
import hashlib
import json
import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional
import zipfile
import xml.etree.ElementTree as ET
from tqdm import tqdm


SUPPORTED_TEXT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".docm",
    ".pptx",
    ".xlsx",
    ".xls",
    ".doc",
}

ARCHIVE_EXTENSIONS = {
    ".zip",
    ".rar",
}

SKIP_EXTENSIONS = {
    ".001",
    ".002",
    ".gif",
}


def stable_document_id(path: Path, root: Path) -> str:
    rel = str(path.relative_to(root)).replace("\\", "/")
    return hashlib.sha1(rel.encode("utf-8", errors="ignore")).hexdigest()[:16]


def write_jsonl(path: Path, item: Dict[str, Any]) -> None:
    with path.open("a", encoding="utf-8") as f:
        f.write(json.dumps(item, ensure_ascii=False) + "\n")


def read_magic_bytes(path: Path, n: int = 16) -> bytes:
    with path.open("rb") as f:
        return f.read(n)


def safe_stat(path: Path) -> Dict[str, Any]:
    try:
        stat = path.stat()
        return {
            "size_bytes": stat.st_size,
            "modified_timestamp": stat.st_mtime,
        }
    except Exception:
        return {
            "size_bytes": None,
            "modified_timestamp": None,
        }


def detect_source_type(path: Path) -> str:
    lower = str(path).lower()

    if "журналы" in lower:
        return "journal"
    if "статьи" in lower:
        return "article"
    if "обзоры" in lower:
        return "review"
    if "доклады" in lower:
        return "presentation_or_report"
    if "материалы конференций" in lower:
        return "conference_material"
    if "источники данных" in lower:
        return "data_source"

    return "unknown"


def parse_pdf(path: Path) -> Dict[str, Any]:
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise RuntimeError("PyMuPDF is not installed. Run: pip install pymupdf") from e

    pages = []

    with fitz.open(path) as doc:
        metadata = doc.metadata or {}

        for i, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            pages.append({
                "page": i,
                "text": text.strip(),
            })

    return {
        "content_type": "pages",
        "metadata": {
            "title": metadata.get("title"),
            "author": metadata.get("author"),
            "subject": metadata.get("subject"),
            "keywords": metadata.get("keywords"),
            "creator": metadata.get("creator"),
            "producer": metadata.get("producer"),
        },
        "pages": pages,
    }

def parse_docm_raw(path: Path) -> Dict[str, Any]:
    """
    Fallback parser for .docm files.
    DOCM is usually an OOXML zip container with word/document.xml.
    We extract text directly from XML without executing macros.
    """
    paragraphs = []

    with zipfile.ZipFile(path, "r") as zf:
        names = zf.namelist()

        if "word/document.xml" not in names:
            raise RuntimeError("DOCM does not contain word/document.xml")

        xml_bytes = zf.read("word/document.xml")

    root = ET.fromstring(xml_bytes)

    namespace = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    }

    for paragraph in root.findall(".//w:p", namespace):
        texts = []
        for node in paragraph.findall(".//w:t", namespace):
            if node.text:
                texts.append(node.text)

        paragraph_text = "".join(texts).strip()
        if paragraph_text:
            paragraphs.append(paragraph_text)

    return {
        "content_type": "document",
        "metadata": {
            "parser": "docm_raw_xml"
        },
        "paragraphs": paragraphs,
        "tables": [],
        "text": "\n".join(paragraphs),
    }
def parse_docx_like(path: Path) -> Dict[str, Any]:
    try:
        import docx
    except ImportError as e:
        raise RuntimeError("python-docx is not installed. Run: pip install python-docx") from e

    document = docx.Document(path)

    paragraphs = []
    for p in document.paragraphs:
        text = p.text.strip()
        if text:
            paragraphs.append(text)

    tables = []
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        tables.append({
            "table_index": table_index,
            "rows": rows,
        })

    return {
        "content_type": "document",
        "metadata": {},
        "paragraphs": paragraphs,
        "tables": tables,
        "text": "\n".join(paragraphs),
    }


def parse_pptx(path: Path) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx") from e

    prs = Presentation(path)
    slides = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)

        slides.append({
            "slide": slide_index,
            "text": "\n".join(texts),
        })

    return {
        "content_type": "slides",
        "metadata": {},
        "slides": slides,
    }


def parse_xlsx(path: Path) -> Dict[str, Any]:
    try:
        import openpyxl
    except ImportError as e:
        raise RuntimeError("openpyxl is not installed. Run: pip install openpyxl") from e

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []

    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            values = []
            for value in row:
                if value is None:
                    values.append("")
                else:
                    values.append(str(value))
            if any(cell.strip() for cell in values):
                rows.append(values)

        sheets.append({
            "sheet_name": ws.title,
            "rows": rows,
            "text": "\n".join(["\t".join(row) for row in rows]),
        })

    return {
        "content_type": "spreadsheet",
        "metadata": {},
        "sheets": sheets,
    }


def parse_xls(path: Path) -> Dict[str, Any]:
    try:
        import pandas as pd
    except ImportError as e:
        raise RuntimeError("pandas is not installed. Run: pip install pandas xlrd") from e
    magic = read_magic_bytes(path)

    if magic.startswith(b"BM"):
        return {
            "content_type": "mislabeled_image",
            "metadata": {
                "detected_format": "bmp",
                "note": "File has .xls extension but BMP signature"
            },
            "sheets": [],
            "text": ""
        }
    sheets_dict = pd.read_excel(path, sheet_name=None, dtype=str, engine="xlrd")
    sheets = []

    for sheet_name, df in sheets_dict.items():
        df = df.fillna("")
        rows = [df.columns.astype(str).tolist()] + df.astype(str).values.tolist()

        sheets.append({
            "sheet_name": str(sheet_name),
            "rows": rows,
            "text": "\n".join(["\t".join(map(str, row)) for row in rows]),
        })

    return {
        "content_type": "spreadsheet",
        "metadata": {},
        "sheets": sheets,
    }


def find_soffice() -> Optional[str]:
    candidates = [
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]

    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved

        if Path(candidate).exists():
            return candidate

    return None


def convert_doc_to_docx(path: Path, temp_dir: Path) -> Path:
    soffice = find_soffice()

    if not soffice:
        raise RuntimeError(
            "Cannot parse .doc: LibreOffice/soffice not found. "
            "Install LibreOffice or convert .doc files manually to .docx."
        )

    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        "docx",
        "--outdir",
        str(temp_dir),
        str(path),
    ]

    result = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        encoding="utf-8",
        errors="ignore",
        timeout=120,
    )

    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed: {result.stderr or result.stdout}"
        )

    converted = temp_dir / f"{path.stem}.docx"

    if not converted.exists():
        candidates = list(temp_dir.glob("*.docx"))
        if candidates:
            return candidates[0]

        raise RuntimeError("LibreOffice finished but converted DOCX was not found.")

    return converted


def parse_doc(path: Path) -> Dict[str, Any]:
    with tempfile.TemporaryDirectory() as tmp:
        temp_dir = Path(tmp)
        docx_path = convert_doc_to_docx(path, temp_dir)
        return parse_docx_like(docx_path)


def parse_file(path: Path) -> Dict[str, Any]:
    ext = path.suffix.lower()

    if ext == ".pdf":
        return parse_pdf(path)

    if ext == ".docx":
        return parse_docx_like(path)

    if ext == ".docm":
        try:
            return parse_docx_like(path)
        except Exception:
            return parse_docm_raw(path)

    if ext == ".pptx":
        return parse_pptx(path)

    if ext == ".xlsx":
        return parse_xlsx(path)

    if ext == ".xls":
        return parse_xls(path)

    if ext == ".doc":
        return parse_doc(path)

    raise ValueError(f"Unsupported extension: {ext}")


def extract_zip(path: Path, output_dir: Path) -> List[Path]:
    extracted_files = []

    target_dir = output_dir / path.stem
    target_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(path, "r") as zf:
        zf.extractall(target_dir)

    for item in target_dir.rglob("*"):
        if item.is_file():
            extracted_files.append(item)

    return extracted_files


def extract_rar(path: Path, output_dir: Path) -> List[Path]:
    try:
        import rarfile
    except ImportError as e:
        raise RuntimeError("rarfile is not installed. Run: pip install rarfile") from e

    extracted_files = []

    target_dir = output_dir / path.stem
    target_dir.mkdir(parents=True, exist_ok=True)

    with rarfile.RarFile(path) as rf:
        rf.extractall(target_dir)

    for item in target_dir.rglob("*"):
        if item.is_file():
            extracted_files.append(item)

    return extracted_files


def collect_files(root: Path) -> List[Path]:
    return [p for p in root.rglob("*") if p.is_file()]


def normalize_text_for_quality_check(parsed: Dict[str, Any]) -> str:
    content_type = parsed.get("content_type")

    if content_type == "pages":
        return "\n".join(page.get("text", "") for page in parsed.get("pages", []))

    if content_type == "document":
        return parsed.get("text", "")

    if content_type == "slides":
        return "\n".join(slide.get("text", "") for slide in parsed.get("slides", []))

    if content_type == "spreadsheet":
        return "\n".join(sheet.get("text", "") for sheet in parsed.get("sheets", []))

    return ""


def process_one_file(
    path: Path,
    root: Path,
    parsed_output_path: Path,
    errors_output_path: Path,
) -> None:
    ext = path.suffix.lower()
    rel_path = str(path.relative_to(root)) if path.is_relative_to(root) else str(path)

    if ext in SKIP_EXTENSIONS:
        write_jsonl(errors_output_path, {
            "path": str(path),
            "relative_path": rel_path,
            "extension": ext,
            "status": "skipped",
            "reason": "Skipped extension",
        })
        return

    if ext not in SUPPORTED_TEXT_EXTENSIONS:
        return

    document_id = stable_document_id(path, root)

    try:
        parsed = parse_file(path)
        full_text = normalize_text_for_quality_check(parsed)

        item = {
            "document_id": document_id,
            "source_path": str(path),
            "relative_path": rel_path,
            "filename": path.name,
            "extension": ext,
            "source_type": detect_source_type(path),
            "file_stat": safe_stat(path),
            "parse_status": "success",
            "text_length": len(full_text),
            "needs_ocr": ext == ".pdf" and len(full_text.strip()) < 200,
            "parsed": parsed,
        }

        write_jsonl(parsed_output_path, item)

    except Exception as e:
        write_jsonl(errors_output_path, {
            "document_id": document_id,
            "source_path": str(path),
            "relative_path": rel_path,
            "filename": path.name,
            "extension": ext,
            "parse_status": "error",
            "error_type": type(e).__name__,
            "error": str(e),
        })


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Parse dataset files into JSONL with extracted raw text."
    )

    parser.add_argument(
        "dataset_path",
        type=str,
        help="Path to dataset folder",
    )

    parser.add_argument(
        "--output-dir",
        type=str,
        default="parsed_output",
        help="Directory for parsed output files",
    )

    parser.add_argument(
        "--extract-archives",
        action="store_true",
        help="Extract ZIP/RAR archives into output directory and parse extracted files too",
    )

    args = parser.parse_args()

    dataset_path = Path(args.dataset_path).resolve()
    output_dir = Path(args.output_dir).resolve()

    if not dataset_path.exists():
        raise FileNotFoundError(f"Dataset path does not exist: {dataset_path}")

    output_dir.mkdir(parents=True, exist_ok=True)

    parsed_output_path = output_dir / "parsed_documents.jsonl"
    errors_output_path = output_dir / "parse_errors.jsonl"
    extracted_archives_dir = output_dir / "extracted_archives"

    if parsed_output_path.exists():
        parsed_output_path.unlink()

    if errors_output_path.exists():
        errors_output_path.unlink()

    all_files = collect_files(dataset_path)

    print(f"Dataset: {dataset_path}")
    print(f"Initial files: {len(all_files)}")
    print(f"Output dir: {output_dir}")

    if args.extract_archives:
        extracted_archives_dir.mkdir(parents=True, exist_ok=True)

        archive_files = [
            p for p in all_files
            if p.suffix.lower() in ARCHIVE_EXTENSIONS
        ]

        print(f"Archives to extract: {len(archive_files)}")

        for archive_path in tqdm(archive_files, desc="Extracting archives"):
            ext = archive_path.suffix.lower()

            try:
                if ext == ".zip":
                    extracted = extract_zip(archive_path, extracted_archives_dir)
                elif ext == ".rar":
                    extracted = extract_rar(archive_path, extracted_archives_dir)
                else:
                    extracted = []

                write_jsonl(errors_output_path, {
                    "source_path": str(archive_path),
                    "extension": ext,
                    "parse_status": "archive_extracted",
                    "extracted_files_count": len(extracted),
                })

            except Exception as e:
                write_jsonl(errors_output_path, {
                    "source_path": str(archive_path),
                    "extension": ext,
                    "parse_status": "archive_extract_error",
                    "error_type": type(e).__name__,
                    "error": str(e),
                })

        extracted_files = collect_files(extracted_archives_dir)
        all_files = all_files + extracted_files

        print(f"Files after archive extraction: {len(all_files)}")

    parse_candidates = [
        p for p in all_files
        if p.suffix.lower() in SUPPORTED_TEXT_EXTENSIONS or p.suffix.lower() in SKIP_EXTENSIONS
    ]

    print(f"Parse candidates: {len(parse_candidates)}")

    for path in tqdm(parse_candidates, desc="Parsing files"):
        process_one_file(
            path=path,
            root=dataset_path,
            parsed_output_path=parsed_output_path,
            errors_output_path=errors_output_path,
        )

    print("\nDone.")
    print(f"Parsed documents: {parsed_output_path}")
    print(f"Errors/logs:       {errors_output_path}")


if __name__ == "__main__":
    main()